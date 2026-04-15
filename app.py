import customtkinter as ctk
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pathlib import Path
import os, shutil, subprocess, tempfile, threading, csv
from typing import List, Dict
from pypdf import PdfReader, PdfWriter

try:
    from PIL import Image, ImageTk
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib import colors as rl_colors
    from reportlab.lib.units import mm
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    import pypdfium2 as pdfium
    PYPDFIUM_AVAILABLE = True
except ImportError:
    PYPDFIUM_AVAILABLE = False

try:
    from pdf2docx import Converter
    PDF2DOCX_AVAILABLE = True
except ImportError:
    PDF2DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# CustomTkinter varsayılan ayarları
ctk.set_appearance_mode("System")  # "Dark", "Light" veya "System"
ctk.set_default_color_theme("blue")  # "blue", "green", "dark-blue"


class PageSelectorDialog:
    def __init__(self, parent, pdf_info: Dict, callback):
        self.pdf_info = pdf_info
        self.callback = callback
        self.dialog = ctk.CTkToplevel(parent)
        self.dialog.title(f"Sayfa Seçimi - {pdf_info['name']}")
        
        # --- GÜVENLİ İKON YÜKLEME ---
        current_dir = os.path.dirname(os.path.abspath(__file__))
        icon_path = os.path.join(current_dir, "pdf.ico")

        if os.path.exists(icon_path):
            try:
                # CustomTkinter'da Toplevel ikonunun ezilmemesi için ufak bir gecikme ekliyoruz
                self.dialog.after(200, lambda: self.dialog.iconbitmap(icon_path))
            except Exception as e:
                print(f"İkon yükleme hatası: {e}")
        else:
            # Hata mesajı dinamikleştirildi
            print(f"Uyarı: {os.path.basename(icon_path)} dosyası bulunamadı.")
        # ----------------------------
        
        self.dialog.geometry("500x650")
        self.dialog.transient(parent)
        self.dialog.grab_set()
        self.create_ui()

    def create_ui(self):
        header = ctk.CTkFrame(self.dialog, height=60, corner_radius=0, fg_color=("gray80", "gray20"))
        header.pack(fill=tk.X)
        ctk.CTkLabel(header, text=f"📄 {self.pdf_info['name']}", font=ctk.CTkFont('Segoe UI', 16, 'bold')).pack(pady=15)
        
        info_frame = ctk.CTkFrame(self.dialog, fg_color="transparent")
        info_frame.pack(fill=tk.X, padx=20, pady=10)
        ctk.CTkLabel(info_frame, text=f"Toplam Sayfa: {self.pdf_info['total_pages']}", font=('Segoe UI', 12)).pack(anchor='w')
        
        quick_frame = ctk.CTkFrame(self.dialog, fg_color="transparent")
        quick_frame.pack(fill=tk.X, padx=20)
        ctk.CTkButton(quick_frame, text="Tümünü Seç", command=self.select_all, width=90).pack(side=tk.LEFT, padx=2)
        ctk.CTkButton(quick_frame, text="Temizle", command=self.deselect_all, width=90, fg_color="#7F8C8D", hover_color="#95A5A6").pack(side=tk.LEFT, padx=2)
        ctk.CTkButton(quick_frame, text="Tekler", command=self.select_odd, width=90).pack(side=tk.LEFT, padx=2)
        ctk.CTkButton(quick_frame, text="Çiftler", command=self.select_even, width=90).pack(side=tk.LEFT, padx=2)
        
        range_frame = ctk.CTkFrame(self.dialog, fg_color="transparent")
        range_frame.pack(fill=tk.X, padx=20, pady=15)
        ctk.CTkLabel(range_frame, text="Sayfa Aralığı (Örn: 1-5, 8, 10-12):", font=('Segoe UI', 12)).pack(anchor='w', pady=(0,5))
        self.range_entry = ctk.CTkEntry(range_frame, font=('Segoe UI', 12))
        self.range_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 5))
        ctk.CTkButton(range_frame, text="Uygula", command=self.apply_range, width=80).pack(side=tk.LEFT)
        
        list_frame = ctk.CTkFrame(self.dialog)
        list_frame.pack(fill=tk.BOTH, expand=True, padx=20, pady=5)
        
        # CTk Scrollbar ile uyumlu klasik Listbox (Binlerce sayfada performans için)
        scrl = ctk.CTkScrollbar(list_frame)
        self.listbox = tk.Listbox(list_frame, yscrollcommand=scrl.set, selectmode=tk.MULTIPLE, 
                                  font=('Segoe UI', 12), bg="#2b2b2b" if ctk.get_appearance_mode()=="Dark" else "#ffffff", 
                                  fg="white" if ctk.get_appearance_mode()=="Dark" else "black", 
                                  selectbackground="#1f538d", highlightthickness=0, borderwidth=0)
        self.listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=2, pady=2)
        scrl.configure(command=self.listbox.yview)
        scrl.pack(side=tk.RIGHT, fill=tk.Y)
        
        for i in range(1, self.pdf_info['total_pages'] + 1):
            self.listbox.insert(tk.END, f"Sayfa {i}")
            if i in self.pdf_info.get('selected_pages', []):
                self.listbox.selection_set(i - 1)
                
        footer = ctk.CTkFrame(self.dialog, fg_color="transparent")
        footer.pack(fill=tk.X, padx=20, pady=20)
        ctk.CTkButton(footer, text="Kaydet", command=self.save_selection, fg_color="#27AE60", hover_color="#229954").pack(side=tk.RIGHT, padx=5)
        ctk.CTkButton(footer, text="İptal", command=self.dialog.destroy, fg_color="#E74C3C", hover_color="#C0392B").pack(side=tk.RIGHT)

    def select_all(self): self.listbox.selection_set(0, tk.END)
    def deselect_all(self): self.listbox.selection_clear(0, tk.END)
    def select_odd(self):
        self.deselect_all()
        for i in range(0, self.pdf_info['total_pages'], 2): self.listbox.selection_set(i)
    def select_even(self):
        self.deselect_all()
        for i in range(1, self.pdf_info['total_pages'], 2): self.listbox.selection_set(i)
    
    def apply_range(self):
        txt = self.range_entry.get().strip()
        if not txt: return
        try:
            pages = set()
            for part in txt.split(','):
                part = part.strip()
                if '-' in part:
                    s, e = map(int, part.split('-'))
                    pages.update(range(s, e + 1))
                else:
                    pages.add(int(part))
            self.deselect_all()
            for p in pages:
                if 1 <= p <= self.pdf_info['total_pages']: self.listbox.selection_set(p - 1)
        except: messagebox.showerror("Hata", "Geçersiz aralık formatı!")

    def save_selection(self):
        sel = [i + 1 for i in self.listbox.curselection()]
        if not sel:
            messagebox.showwarning("Uyarı", "En az bir sayfa seçmelisiniz!")
            return
        self.pdf_info['selected_pages'] = sel
        self.callback()
        self.dialog.destroy()


class PDFToolboxApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF Toolbox")
        self.root.geometry("1000x850")
        
        # --- GÜVENLİ İKON YÜKLEME ---
        # Dosyanın tam yolunu oluşturuyoruz
        current_dir = os.path.dirname(os.path.abspath(__file__))
        icon_path = os.path.join(current_dir, "pdf.ico")

        if os.path.exists(icon_path):
            try:
                # Bazı Windows sürümlerinde wm_iconbitmap daha kararlı çalışır
                self.root.after(200, lambda: self.root.wm_iconbitmap(icon_path))
            except Exception as e:
                print(f"İkon yükleme hatası: {e}")
        else:
            print("Uyarı: pdf.ico dosyası bulunamadı.")
        # ----------------------------
        
        
        self.pdf_files: List[Dict] = []
        self.split_data = {'path': None, 'name': "", 'total_pages': 0, 'selected_pages': []}
        self.edit_data = {'path': None, 'name': "", 'total_pages': 0, 'current_page': 1, 'preview_image': None, 
                          'edits': [], 'selected_tool': None, 'temp_params': {}, 'scale_factor': 1.0, 
                          'canvas_items': [], 'drag_start': None}
        self.export_data = {'path': None, 'name': '', 'total_pages': 0, 'selected_pages': []}
        
        self.setup_styles()
        self.create_ui()

    def setup_styles(self):
        # Treeview'i CustomTkinter temasıyla uyumlu hale getirmek için
        style = ttk.Style()
        style.theme_use('default')
        bg_color = "#2b2b2b" if ctk.get_appearance_mode() == "Dark" else "#ffffff"
        fg_color = "white" if ctk.get_appearance_mode() == "Dark" else "black"
        head_bg = "#1f538d"
        
        style.configure("Treeview", background=bg_color, foreground=fg_color, fieldbackground=bg_color, rowheight=35, borderwidth=0, font=('Segoe UI', 11))
        style.configure("Treeview.Heading", background=head_bg, foreground="white", font=('Segoe UI', 11, 'bold'), relief="flat")
        style.map("Treeview", background=[('selected', head_bg)])

    def create_ui(self):
        header = ctk.CTkFrame(self.root, height=70, corner_radius=0, fg_color=("#34495E", "#1a1a1a"))
        header.pack(fill=tk.X)
        ctk.CTkLabel(header, text="📝 PDF Toolbox", font=ctk.CTkFont('Segoe UI', 24, 'bold'), text_color="white").pack(pady=8)

        self.notebook = ctk.CTkTabview(self.root)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)

        self.merge_tab = self.notebook.add("PDF Birleştir")
        self.split_tab = self.notebook.add("PDF Böl")
        self.convert_tab = self.notebook.add("Dosya → PDF")
        self.edit_tab = self.notebook.add("PDF Düzenle")
        self.export_tab = self.notebook.add("PDF → Diğer")

        self.setup_merge_ui()
        self.setup_split_ui()
        self.setup_convert_ui()
        self.setup_edit_ui()
        self.setup_export_ui()

        self.status_label = ctk.CTkLabel(self.root, text="📌 Hazır", font=('Segoe UI', 12), anchor='w')
        self.status_label.pack(fill=tk.X, padx=20, pady=(0, 2))
        self.status_label =ctk.CTkLabel(self.root, text="Geliştirici: Emrullah ALKAÇ", font=ctk.CTkFont("Arial", 12, "bold"), text_color="#2980B9").pack(anchor="center", padx=2, pady=(2,2))
        

    # ================= BİRLEŞTİRME SEKMESİ =================
    def setup_merge_ui(self):
        content = ctk.CTkFrame(self.merge_tab, fg_color="transparent")
        content.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        btn_frame = ctk.CTkFrame(content, width=180, fg_color="transparent")
        btn_frame.pack(side=tk.LEFT, fill=tk.Y, padx=(0, 15))
        
        btns = [
            ("➕ PDF Ekle", self.add_pdf_files, "#27AE60", "#229954"),
            ("➖ Seçileni Çıkar", self.remove_selected, "#E74C3C", "#C0392B"),
            ("⬆ Yukarı Taşı", self.move_up, "#3498DB", "#2980B9"),
            ("⬇ Aşağı Taşı", self.move_down, "#3498DB", "#2980B9"),
            ("📑 Sayfa Seç", self.select_pages, "#9B59B6", "#8E44AD"),
            ("🗑 Temizle", self.clear_all, "#538F93", "#46959B")
        ]
        
        for text, cmd, bg, hov in btns:
            ctk.CTkButton(btn_frame, text=text, command=cmd, fg_color=bg, hover_color=hov, width=160, height=40).pack(pady=5)
            
        ctk.CTkFrame(btn_frame, fg_color="transparent", height=40).pack()
        ctk.CTkButton(btn_frame, text="🔗 BİRLEŞTİR", command=self.merge_pdfs, fg_color="#E67E22", hover_color="#D35400", width=160, height=40, font=('Segoe UI', 14, 'bold')).pack(side=tk.BOTTOM)
        
        tree_frame = ctk.CTkFrame(content)
        tree_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        self.tree = ttk.Treeview(tree_frame, columns=('Dosya', 'Sayfa', 'Seçili'), show='headings')
        self.tree.heading('Dosya', text='Dosya Adı')
        self.tree.heading('Sayfa', text='Toplam Sayfa')
        self.tree.heading('Seçili', text='Seçili Aralık')
        self.tree.column('Dosya', width=300)
        self.tree.column('Sayfa', width=100, anchor='center')
        self.tree.column('Seçili', width=180, anchor='center')
        
        scrl = ctk.CTkScrollbar(tree_frame, orientation="vertical", command=self.tree.yview)
        self.tree.configure(yscrollcommand=scrl.set)
        self.tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(2,0), pady=2)
        scrl.pack(side=tk.RIGHT, fill=tk.Y, padx=2, pady=2)

    def add_pdf_files(self):
        files = filedialog.askopenfilenames(filetypes=[("PDF Dosyaları", "*.pdf")])
        for f in files:
            try:
                reader = PdfReader(f)
                info = {'path': f, 'name': Path(f).name, 'total_pages': len(reader.pages), 'selected_pages': list(range(1, len(reader.pages) + 1))}
                self.pdf_files.append(info)
            except: continue
        self.refresh_tree()

    def refresh_tree(self):
        for item in self.tree.get_children(): self.tree.delete(item)
        for f in self.pdf_files:
            p_range = "Tümü" if len(f['selected_pages']) == f['total_pages'] else self.format_page_ranges(f['selected_pages'])
            self.tree.insert('', 'end', values=(f['name'], f['total_pages'], p_range))

    def remove_selected(self):
        sel = self.tree.selection()
        if not sel: return
        del self.pdf_files[self.tree.index(sel[0])]
        self.refresh_tree()

    def move_up(self):
        sel = self.tree.selection()
        if not sel: return
        idx = self.tree.index(sel[0])
        if idx > 0:
            self.pdf_files[idx], self.pdf_files[idx-1] = self.pdf_files[idx-1], self.pdf_files[idx]
            self.refresh_tree()
            self.tree.selection_set(self.tree.get_children()[idx-1])

    def move_down(self):
        sel = self.tree.selection()
        if not sel: return
        idx = self.tree.index(sel[0])
        if idx < len(self.pdf_files)-1:
            self.pdf_files[idx], self.pdf_files[idx+1] = self.pdf_files[idx+1], self.pdf_files[idx]
            self.refresh_tree()
            self.tree.selection_set(self.tree.get_children()[idx+1])

    def select_pages(self):
        sel = self.tree.selection()
        if not sel: return messagebox.showwarning("Uyarı", "Lütfen bir dosya seçin!")
        idx = self.tree.index(sel[0])
        PageSelectorDialog(self.root, self.pdf_files[idx], self.refresh_tree)

    def format_page_ranges(self, pages: List[int]) -> str:
        if not pages: return "Hiçbiri"
        pages = sorted(pages)
        ranges, start, end = [], pages[0], pages[0]
        for i in range(1, len(pages)):
            if pages[i] == end + 1: end = pages[i]
            else:
                ranges.append(f"{start}-{end}" if start != end else str(start))
                start = end = pages[i]
        ranges.append(f"{start}-{end}" if start != end else str(start))
        res = ", ".join(ranges)
        return res if len(res) < 40 else res[:37] + "..."

    def merge_pdfs(self):
        if not self.pdf_files: return
        out = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Dosyaları", "*.pdf")])
        if not out: return
        try:
            writer = PdfWriter()
            for info in self.pdf_files:
                reader = PdfReader(info['path'])
                for p in sorted(info['selected_pages']): writer.add_page(reader.pages[p-1])
            with open(out, "wb") as f: writer.write(f)
            messagebox.showinfo("Başarılı", "PDF başarıyla oluşturuldu.")
        except Exception as e: messagebox.showerror("Hata", str(e))

    def clear_all(self):
        self.pdf_files.clear()
        self.refresh_tree()

    # ================= BÖLME SEKMESİ =================
    def setup_split_ui(self):
        container = ctk.CTkFrame(self.split_tab, fg_color="transparent")
        container.pack(pady=40, padx=50, fill=tk.BOTH)
        
        ctk.CTkLabel(container, text="Bölünecek PDF Dosyası", font=ctk.CTkFont('Segoe UI', 16, 'bold')).pack(anchor='w')
        
        f_frame = ctk.CTkFrame(container, fg_color="transparent")
        f_frame.pack(fill=tk.X, pady=10)
        self.split_file_label = ctk.CTkLabel(f_frame, text="Lütfen bir dosya seçin...", corner_radius=6, fg_color=("gray85", "gray20"), anchor="w", padx=10, height=40)
        self.split_file_label.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        ctk.CTkButton(f_frame, text="Dosya Seç", command=self.select_split_pdf, width=120, height=40).pack(side=tk.RIGHT)
        
        selection_frame = ctk.CTkFrame(container, fg_color="transparent")
        selection_frame.pack(fill=tk.X, pady=10)
        self.split_info = ctk.CTkLabel(selection_frame, text="Seçili Sayfalar: -", font=('Segoe UI', 12))
        self.split_info.pack(side=tk.LEFT)
        ctk.CTkButton(selection_frame, text="Sayfaları Düzenle", command=self.open_split_page_selector, fg_color="#9B59B6", hover_color="#8E44AD", width=150, height=35).pack(side=tk.RIGHT)
        
        settings_box = ctk.CTkFrame(container)
        settings_box.pack(fill=tk.X, pady=20)
        ctk.CTkLabel(settings_box, text=" Bölme Ayarları ", font=ctk.CTkFont('Segoe UI', 14, 'bold')).pack(anchor="w", padx=15, pady=(15, 0))
        
        inner_set = ctk.CTkFrame(settings_box, fg_color="transparent")
        inner_set.pack(fill=tk.X, padx=15, pady=15)
        ctk.CTkLabel(inner_set, text="Her yeni dosyada kaç sayfa olsun?", font=('Segoe UI', 13)).pack(side=tk.LEFT)
        self.pages_per_file_entry = ctk.CTkEntry(inner_set, font=('Segoe UI', 14), width=60, justify='center')
        self.pages_per_file_entry.insert(0, "1")
        self.pages_per_file_entry.pack(side=tk.LEFT, padx=15)
        ctk.CTkLabel(inner_set, text="(Seçtiğiniz sayfalar bu düzende parçalanacaktır)", text_color="gray").pack(side=tk.LEFT)
        
        self.split_progress = ctk.CTkProgressBar(container)
        self.split_progress.set(0)
        self.split_progress.pack(fill=tk.X, pady=20)
        
        ctk.CTkButton(container, text="✂️ PDF'İ BÖL VE KAYDET", command=self.execute_split, fg_color="#27AE60", hover_color="#229954", font=('Segoe UI', 16, 'bold'), height=50).pack(pady=8)

    def select_split_pdf(self):
        path = filedialog.askopenfilename(filetypes=[("PDF Dosyaları", "*.pdf")])
        if path:
            reader = PdfReader(path)
            total = len(reader.pages)
            self.split_data = {'path': path, 'name': Path(path).name, 'total_pages': total, 'selected_pages': list(range(1, total + 1))}
            self.split_file_label.configure(text=self.split_data['name'])
            self.update_split_ui_info()

    def open_split_page_selector(self):
        if not self.split_data['path']: return messagebox.showwarning("Uyarı", "Lütfen önce bir PDF dosyası seçin!")
        PageSelectorDialog(self.root, self.split_data, self.update_split_ui_info)

    def update_split_ui_info(self):
        range_text = self.format_page_ranges(self.split_data['selected_pages'])
        self.split_info.configure(text=f"Seçili Sayfalar ({len(self.split_data['selected_pages'])}): {range_text}")

    def execute_split(self):
        if not self.split_data['path']: return messagebox.showwarning("Uyarı", "Dosya seçilmedi!")
        try:
            step = int(self.pages_per_file_entry.get())
            if step <= 0: raise ValueError
            output_dir = filedialog.askdirectory()
            if not output_dir: return
            
            reader = PdfReader(self.split_data['path'])
            selected_pages = sorted(self.split_data['selected_pages'])
            total_selected = len(selected_pages)
            self.split_progress.set(0)
            
            progress_counter = 0
            for i in range(0, total_selected, step):
                writer = PdfWriter()
                current_chunk = selected_pages[i : i + step]
                for page_num in current_chunk:
                    writer.add_page(reader.pages[page_num - 1])
                    progress_counter += 1
                    self.split_progress.set(progress_counter / total_selected)
                    self.root.update_idletasks()
                file_name = f"{Path(self.split_data['path']).stem}_part_{i//step + 1}.pdf"
                with open(os.path.join(output_dir, file_name), "wb") as f: writer.write(f)
                
            self.split_progress.set(0)
            messagebox.showinfo("Başarılı", f"Bölme işlemi tamamlandı.\n{total_selected} sayfa, {(total_selected + step - 1)//step} dosyaya bölündü.")
        except ValueError: messagebox.showerror("Hata", "Lütfen geçerli bir sayfa sayısı girin.")
        except Exception as e: messagebox.showerror("Hata", str(e))

    # ================= DOSYA → PDF DÖNÜŞTÜRME =================
    def setup_convert_ui(self):
        container = ctk.CTkFrame(self.convert_tab, fg_color="transparent")
        container.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        info = ctk.CTkFrame(container, fg_color=("#D5E8D4", "#2c4a2c"), corner_radius=6)
        info.pack(fill=tk.X, pady=(0, 10))
        ctk.CTkLabel(info, text="📂 Desteklenen formatlar: DOC · DOCX · ODT · PNG · JPG / JPEG → PDF", font=('Segoe UI', 13), text_color=("#1A5E1A", "#a3d9a3")).pack(pady=8)
        
        mid = ctk.CTkFrame(container, fg_color="transparent")
        mid.pack(fill=tk.BOTH, expand=True)
        
        btn_frame = ctk.CTkFrame(mid, width=180, fg_color="transparent")
        btn_frame.pack(side=tk.LEFT, fill=tk.Y, padx=(0, 15))
        
        btn_list = [
            ("➕ Dosya Ekle", self.conv_add_files, "#27AE60", "#229954"),
            ("➖ Seçileni Çıkar", self.conv_remove_selected, "#E74C3C", "#C0392B"),
            ("⬆ Yukarı Taşı", self.conv_move_up, "#3498DB", "#2980B9"),
            ("⬇ Aşağı Taşı", self.conv_move_down, "#3498DB", "#2980B9"),
            ("🗑 Listeyi Temizle", self.conv_clear, "#409AA1", "#49878B"),
            ("💾 Tek Kaydet", self.conv_export_single, "#8E44AD", "#6C3483")
        ]
        for txt, cmd, bg, hov in btn_list:
            ctk.CTkButton(btn_frame, text=txt, command=cmd, fg_color=bg, hover_color=hov, width=160, height=40).pack(pady=5)
            
        ctk.CTkFrame(btn_frame, fg_color="transparent", height=20).pack()
        ctk.CTkButton(btn_frame, text="🔗 BİRLEŞTİR\n& DÖNÜŞTÜR", command=self.conv_export_merged, fg_color="#E67E22", hover_color="#D35400", width=160, height=40, font=('Segoe UI', 13, 'bold')).pack(side=tk.BOTTOM, pady=5)
        
        list_frame = ctk.CTkFrame(mid)
        list_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        self.conv_tree = ttk.Treeview(list_frame, columns=('Dosya', 'Tür', 'Durum'), show='headings')
        self.conv_tree.heading('Dosya', text='Dosya Adı')
        self.conv_tree.heading('Tür', text='Tür')
        self.conv_tree.heading('Durum', text='Durum')
        self.conv_tree.column('Dosya', width=380)
        self.conv_tree.column('Tür', width=80, anchor='center')
        self.conv_tree.column('Durum', width=120, anchor='center')
        
        scrl = ctk.CTkScrollbar(list_frame, orientation="vertical", command=self.conv_tree.yview)
        self.conv_tree.configure(yscrollcommand=scrl.set)
        self.conv_tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(2,0), pady=2)
        scrl.pack(side=tk.RIGHT, fill=tk.Y, padx=2, pady=2)
        
        bottom = ctk.CTkFrame(container, fg_color="transparent")
        bottom.pack(fill=tk.X, pady=(10, 0))
        self.conv_progress = ctk.CTkProgressBar(bottom)
        self.conv_progress.set(0)
        self.conv_progress.pack(fill=tk.X)
        self.conv_status = ctk.CTkLabel(bottom, text="Hazır.", font=('Segoe UI', 11), text_color="gray", anchor='w')
        self.conv_status.pack(anchor='w', pady=2)
        
        self.conv_files: List[Dict] = []

    def conv_add_files(self):
        filetypes = [("Desteklenen Dosyalar", "*.doc *.docx *.odt *.png *.jpg *.jpeg"), ("Tüm Dosyalar", "*.*")]
        files = filedialog.askopenfilenames(filetypes=filetypes)
        for f in files: self.conv_files.append({'path': f, 'name': Path(f).name, 'status': 'Bekliyor'})
        self.conv_refresh_tree()

    def conv_refresh_tree(self):
        for item in self.conv_tree.get_children(): self.conv_tree.delete(item)
        for f in self.conv_files: self.conv_tree.insert('', 'end', values=(f['name'], Path(f['path']).suffix.upper().lstrip('.'), f.get('status', '—')))

    def conv_remove_selected(self):
        sel = self.conv_tree.selection()
        if not sel: return
        del self.conv_files[self.conv_tree.index(sel[0])]
        self.conv_refresh_tree()

    def conv_move_up(self):
        sel = self.conv_tree.selection()
        if not sel: return
        idx = self.conv_tree.index(sel[0])
        if idx > 0:
            self.conv_files[idx], self.conv_files[idx-1] = self.conv_files[idx-1], self.conv_files[idx]
            self.conv_refresh_tree()
            self.conv_tree.selection_set(self.conv_tree.get_children()[idx-1])

    def conv_move_down(self):
        sel = self.conv_tree.selection()
        if not sel: return
        idx = self.conv_tree.index(sel[0])
        if idx < len(self.conv_files)-1:
            self.conv_files[idx], self.conv_files[idx+1] = self.conv_files[idx+1], self.conv_files[idx]
            self.conv_refresh_tree()
            self.conv_tree.selection_set(self.conv_tree.get_children()[idx+1])

    def conv_clear(self):
        self.conv_files.clear()
        self.conv_refresh_tree()

    def _convert_single_to_pdf(self, file_info: Dict, tmp_dir: str) -> str:
        src = file_info['path']
        ext = Path(src).suffix.lower()
        out_pdf = os.path.join(tmp_dir, Path(src).stem + ".pdf")
        if ext in ('.doc', '.docx', '.odt'):
            local_src = os.path.join(tmp_dir, "input" + ext)
            shutil.copy2(src, local_src)
            profile_dir = os.path.join(tmp_dir, "lo_profile")
            os.makedirs(profile_dir, exist_ok=True)
            result = subprocess.run([shutil.which('libreoffice') or shutil.which('soffice') or r'C:\Program Files\LibreOffice\program\soffice.exe',
                                     '--headless', '--norestore', '--nofirststartwizard',
                                     f'-env:UserInstallation={Path(profile_dir).as_uri()}', '--convert-to', 'pdf', '--outdir', tmp_dir, local_src],
                                    capture_output=True, timeout=120)
            if result.returncode != 0: raise RuntimeError(f"LibreOffice hatası: {result.stderr.decode('utf-8', errors='replace')}")
            out_pdf = os.path.join(tmp_dir, [f for f in os.listdir(tmp_dir) if f.endswith('.pdf')][0])
        elif ext in ('.png', '.jpg', '.jpeg'):
            if not PIL_AVAILABLE: raise RuntimeError("Pillow kurulu değil: pip install Pillow")
            img = Image.open(src).convert('RGB')
            img.save(out_pdf, 'PDF', resolution=150)
        else: raise RuntimeError(f"Desteklenmeyen tür: {ext}")
        return out_pdf

    def conv_export_single(self):
        if not self.conv_files: return messagebox.showwarning("Uyarı", "Dosya ekleyin!")
        out_dir = filedialog.askdirectory()
        if not out_dir: return
        def run():
            total = len(self.conv_files)
            self.root.after(0, lambda: self.conv_progress.set(0))
            success = 0
            for i, f in enumerate(self.conv_files):
                self.root.after(0, lambda text=f"Dönüştürülüyor: {f['name']}": self.conv_status.configure(text=text))
                self.root.update_idletasks()
                with tempfile.TemporaryDirectory() as tmp:
                    try:
                        pdf_path = self._convert_single_to_pdf(f, tmp)
                        dest = os.path.join(out_dir, f"{Path(f['name']).stem}.pdf")
                        if os.path.exists(dest): dest = os.path.join(out_dir, f"{Path(f['name']).stem}_{i+1}.pdf")
                        shutil.copy2(pdf_path, dest)
                        self.conv_files[i]['status'] = '✅ Tamam'
                        success += 1
                    except Exception as e:
                        self.conv_files[i]['status'] = '❌ Hata'
                        self.root.after(0, lambda text=f"Hata: {e}": self.conv_status.configure(text=text))
                self.root.after(0, lambda v=(i+1)/total: self.conv_progress.set(v))
                self.root.after(0, self.conv_refresh_tree)
                self.root.update_idletasks()
            self.root.after(0, lambda: self.conv_status.configure(text=f"Tamamlandı. {success}/{total} dosya."))
            messagebox.showinfo("Bitti", f"{success} dosya kaydedildi: {out_dir}")
        threading.Thread(target=run, daemon=True).start()

    def conv_export_merged(self):
        if not self.conv_files: return messagebox.showwarning("Uyarı", "Dosya ekleyin!")
        out_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Dosyaları", "*.pdf")])
        if not out_path: return
        def run():
            total = len(self.conv_files)
            self.root.after(0, lambda: self.conv_progress.set(0))
            writer = PdfWriter()
            success = 0
            for i, f in enumerate(self.conv_files):
                self.root.after(0, lambda text=f"Dönüştürülüyor: {f['name']}": self.conv_status.configure(text=text))
                with tempfile.TemporaryDirectory() as tmp:
                    try:
                        pdf_path = self._convert_single_to_pdf(f, tmp)
                        for page in PdfReader(pdf_path).pages: writer.add_page(page)
                        self.conv_files[i]['status'] = '✅ Eklendi'
                        success += 1
                    except Exception as e:
                        self.conv_files[i]['status'] = '❌ Hata'
                self.root.after(0, lambda v=(i+1)/total: self.conv_progress.set(v))
                self.root.after(0, self.conv_refresh_tree)
            if writer.pages:
                with open(out_path, 'wb') as fp: writer.write(fp)
                self.root.after(0, lambda: self.conv_status.configure(text=f"{success} dosya birleştirildi."))
                messagebox.showinfo("Başarılı", out_path)
            else:
                self.root.after(0, lambda: self.conv_status.configure(text="Hata: Sayfa eklenemedi."))
        threading.Thread(target=run, daemon=True).start()

    # ================= PDF DÜZENLEME SEKMESİ =================
    def setup_edit_ui(self):
        if not REPORTLAB_AVAILABLE:
            ctk.CTkLabel(self.edit_tab, text="⚠️ PDF Düzenleme için reportlab gerekli\npip install reportlab", font=('Segoe UI', 16), text_color="#E74C3C").pack(expand=True)
            return
            
        container = ctk.CTkFrame(self.edit_tab, fg_color="transparent")
        container.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        
        tools_panel = ctk.CTkFrame(container, width=200)
        tools_panel.pack(side=tk.LEFT, fill=tk.Y, padx=(0, 10))
        tools_panel.pack_propagate(False)
        
        ctk.CTkLabel(tools_panel, text="🎨 Araçlar", font=ctk.CTkFont('Segoe UI', 14, 'bold'), fg_color=("#34495E", "#1a1a1a"), text_color="white", corner_radius=6).pack(fill=tk.X, pady=5, padx=5)
        
        for txt, cmd in [("📂 PDF Aç", self.edit_open_pdf), ("💾 Kaydet", self.edit_save_pdf)]:
            ctk.CTkButton(tools_panel, text=txt, command=cmd).pack(fill=tk.X, padx=10, pady=4)
            
        ctk.CTkFrame(tools_panel, height=2).pack(fill=tk.X, pady=5, padx=10)
        ctk.CTkLabel(tools_panel, text="Metin & Şekil", font=('Segoe UI', 12, 'bold')).pack(fill=tk.X, padx=10, pady=(0, 5))
        
        self.edit_tool_buttons = {}
        for tool_id, text in [("text", "✏️ Metin"), ("rectangle", "⬜ Dikdörtgen"), ("circle", "⭕ Daire"), ("line", "📏 Çizgi")]:
            btn = ctk.CTkButton(tools_panel, text=text, fg_color="transparent", border_width=1, border_color="gray50", text_color=("black", "white"), command=lambda t=tool_id: self.edit_select_tool(t))
            btn.pack(fill=tk.X, padx=10, pady=3)
            self.edit_tool_buttons[tool_id] = btn
            
        for txt, cmd in [("🖼️ Resim Ekle", self.edit_add_image), ("💧 Filigran", self.edit_add_watermark), ("🔄 Döndür", self.edit_rotate_page), ("❌ Sayfa Sil", self.edit_delete_page)]:
            ctk.CTkButton(tools_panel, text=txt, command=cmd, fg_color="transparent", border_width=1, border_color="gray50", text_color=("black", "white")).pack(fill=tk.X, padx=10, pady=3)
            
        ctk.CTkFrame(tools_panel, height=2).pack(fill=tk.X, pady=10, padx=10)
        
        self.edit_color_var = ctk.StringVar(value="black")
        ctk.CTkComboBox(tools_panel, variable=self.edit_color_var, values=["black","red","blue","green","white"], state="readonly").pack(padx=10, pady=4, fill=tk.X)
        
        self.edit_size_var = ctk.StringVar(value="12")
        ctk.CTkComboBox(tools_panel, variable=self.edit_size_var, values=[str(i) for i in range(6, 74, 2)], state="readonly").pack(padx=10, pady=4, fill=tk.X)
        
        self.edit_fill_var = ctk.BooleanVar(value=False)
        ctk.CTkCheckBox(tools_panel, text="Dolgulu Şekil", variable=self.edit_fill_var).pack(anchor='w', padx=10, pady=5)
        
        ctk.CTkFrame(tools_panel, height=2).pack(fill=tk.X, pady=4, padx=10)
        ctk.CTkButton(tools_panel, text="📋 Liste", command=self.edit_show_edits).pack(fill=tk.X, padx=10, pady=2)
        ctk.CTkButton(tools_panel, text="🗑️ Temizle", command=self.edit_clear_all, fg_color="#E74C3C", hover_color="#C0392B").pack(fill=tk.X, padx=10, pady=2)
        
        right_panel = ctk.CTkFrame(container, fg_color="transparent")
        right_panel.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        info_frame = ctk.CTkFrame(right_panel)
        info_frame.pack(fill=tk.X, pady=(0, 10))
        self.edit_file_label = ctk.CTkLabel(info_frame, text="📄 Dosya seçilmedi", font=ctk.CTkFont('Segoe UI', 13, 'bold'), anchor='w')
        self.edit_file_label.pack(fill=tk.X, padx=10, pady=5)
        
        nav_frame = ctk.CTkFrame(info_frame, fg_color="transparent")
        nav_frame.pack(fill=tk.X, padx=10, pady=5)
        ctk.CTkButton(nav_frame, text="◀", command=self.edit_prev_page, width=40).pack(side=tk.LEFT, padx=2)
        self.edit_page_label = ctk.CTkLabel(nav_frame, text="Sayfa: -", font=('Segoe UI', 12))
        self.edit_page_label.pack(side=tk.LEFT, padx=20)
        ctk.CTkButton(nav_frame, text="▶", command=self.edit_next_page, width=40).pack(side=tk.LEFT, padx=2)
        ctk.CTkButton(nav_frame, text="🔄 Yenile", command=self.edit_refresh_preview, width=80).pack(side=tk.RIGHT, padx=2)
        
        # Canvas, CTk arayüzü ile uyumlu olması için tk.Canvas olarak bırakılmıştır, arka plan rengi ayarlandı.
        canvas_frame = ctk.CTkFrame(right_panel)
        canvas_frame.pack(fill=tk.BOTH, expand=True)
        scrollbar_v = ctk.CTkScrollbar(canvas_frame, orientation=tk.VERTICAL)
        scrollbar_h = ctk.CTkScrollbar(canvas_frame, orientation=tk.HORIZONTAL)
        
        canvas_bg = "#e0e0e0" if ctk.get_appearance_mode() == "Light" else "#202020"
        self.edit_canvas = tk.Canvas(canvas_frame, bg=canvas_bg, highlightthickness=0, yscrollcommand=scrollbar_v.set, xscrollcommand=scrollbar_h.set)
        
        scrollbar_v.configure(command=self.edit_canvas.yview)
        scrollbar_h.configure(command=self.edit_canvas.xview)
        scrollbar_v.pack(side=tk.RIGHT, fill=tk.Y)
        scrollbar_h.pack(side=tk.BOTTOM, fill=tk.X)
        self.edit_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        
        self.edit_canvas.bind("<Button-1>", self.edit_canvas_click)
        self.edit_canvas.bind("<B1-Motion>", self.edit_canvas_drag)
        self.edit_canvas.bind("<ButtonRelease-1>", self.edit_canvas_release)
        self.edit_canvas.bind("<Motion>", self.edit_canvas_motion)
        
        status_frame = ctk.CTkFrame(right_panel, fg_color="transparent")
        status_frame.pack(fill=tk.X, pady=(5,0))
        self.edit_tool_label = ctk.CTkLabel(status_frame, text="🖱️ Araç: Yok", font=('Segoe UI', 11, 'bold'), anchor='w')
        self.edit_tool_label.pack(side=tk.LEFT)
        self.edit_status = ctk.CTkLabel(status_frame, text="Hazır", font=('Segoe UI', 11), text_color="gray", anchor='w')
        self.edit_status.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=10)

    def edit_open_pdf(self):
        file_path = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if not file_path: return
        try:
            reader = PdfReader(file_path)
            self.edit_data.update({'path': file_path, 'name': Path(file_path).name, 'total_pages': len(reader.pages), 'current_page': 1, 'edits': []})
            self.edit_file_label.configure(text=f"📄 {self.edit_data['name']}")
            self.edit_page_label.configure(text=f"Sayfa: 1 / {self.edit_data['total_pages']}")
            self.edit_status.configure(text=f"PDF yüklendi: {self.edit_data['total_pages']} sayfa")
            self.edit_refresh_preview()
        except Exception as e: messagebox.showerror("Hata", str(e))

    def edit_refresh_preview(self):
        if not self.edit_data['path'] or not PYPDFIUM_AVAILABLE: return
        try:
            self.edit_status.configure(text="Önizleme oluşturuluyor...")
            self.root.update_idletasks()
            pdf = pdfium.PdfDocument(self.edit_data['path'])
            page = pdf[self.edit_data['current_page'] - 1]
            pil_image = page.render(scale=2.0).to_pil()
            pdf_w = float(PdfReader(self.edit_data['path']).pages[self.edit_data['current_page'] - 1].mediabox.width)
            pdf_h = float(PdfReader(self.edit_data['path']).pages[self.edit_data['current_page'] - 1].mediabox.height)
            scale = min(700 / pil_image.size[0], 900 / pil_image.size[1], 1.0)
            if scale < 1.0: pil_image = pil_image.resize((int(pil_image.size[0]*scale), int(pil_image.size[1]*scale)), Image.LANCZOS)
            self.edit_data['scale_factor'] = pdf_w / pil_image.size[0]
            self.edit_data['pdf_height'] = pdf_h
            self.edit_data['preview_image'] = ImageTk.PhotoImage(pil_image)
            self.edit_canvas.delete("all")
            self.edit_canvas.create_image(0, 0, anchor=tk.NW, image=self.edit_data['preview_image'])
            self.edit_redraw_edits()
            self.edit_canvas.config(scrollregion=self.edit_canvas.bbox(tk.ALL))
            self.edit_status.configure(text=f"Önizleme hazır - Sayfa {self.edit_data['current_page']}")
        except Exception as e: messagebox.showerror("Hata", str(e))

    def edit_redraw_edits(self):
        for item in self.edit_data.get('canvas_items', []):
            try: self.edit_canvas.delete(item)
            except: pass
        self.edit_data['canvas_items'] = []
        if not self.edit_data['path']: return
        page_edits = [e for e in self.edit_data['edits'] if e['page'] == self.edit_data['current_page']]
        scale = self.edit_data.get('scale_factor', 1.0)
        pdf_h = self.edit_data.get('pdf_height', 842)
        for edit in page_edits:
            p = edit['params']
            if edit['type'] == 'text':
                self.edit_data['canvas_items'].append(self.edit_canvas.create_text(p['x']/scale, (pdf_h-p['y'])/scale, text=p['text'], fill=p['color'], font=('Helvetica', max(8, int(p['size']/scale))), anchor='nw'))
            elif edit['type'] == 'rectangle':
                self.edit_data['canvas_items'].append(self.edit_canvas.create_rectangle(p['x']/scale, (pdf_h-p['y'])/scale, (p['x']+p['width'])/scale, (pdf_h-p['y']-p['height'])/scale, fill=p['fill'] if p['fill']!='none' else '', outline=p['stroke'] if p['stroke']!='none' else ''))
            elif edit['type'] == 'circle':
                self.edit_data['canvas_items'].append(self.edit_canvas.create_oval((p['x']-p['radius'])/scale, (pdf_h-p['y']-p['radius'])/scale, (p['x']+p['radius'])/scale, (pdf_h-p['y']+p['radius'])/scale, fill=p['fill'] if p['fill']!='none' else '', outline=p['stroke'] if p['stroke']!='none' else ''))
            elif edit['type'] == 'line':
                self.edit_data['canvas_items'].append(self.edit_canvas.create_line(p['x1']/scale, (pdf_h-p['y1'])/scale, p['x2']/scale, (pdf_h-p['y2'])/scale, fill=p['color'], width=p['width']))

    def edit_select_tool(self, tool):
        if self.edit_data['selected_tool']:
            old_btn = self.edit_tool_buttons.get(self.edit_data['selected_tool'])
            if old_btn: old_btn.configure(fg_color="transparent")
        if self.edit_data['selected_tool'] == tool:
            self.edit_data['selected_tool'] = None
            self.edit_tool_label.configure(text="🖱️ Araç: Yok")
            self.edit_canvas.config(cursor="")
        else:
            self.edit_data['selected_tool'] = tool
            self.edit_tool_buttons[tool].configure(fg_color="#4CAF50")
            self.edit_tool_label.configure(text=f"🖱️ {tool}")
            self.edit_canvas.config(cursor="crosshair")

    def edit_canvas_click(self, event):
        if not self.edit_data['path'] or not self.edit_data['selected_tool']: return
        cx, cy = self.edit_canvas.canvasx(event.x), self.edit_canvas.canvasy(event.y)
        if self.edit_data['selected_tool'] == 'text': self.edit_add_text_at_position(cx, cy)
        else: self.edit_data['drag_start'] = (cx, cy)

    def edit_canvas_drag(self, event):
        if not self.edit_data.get('drag_start'): return
        cx, cy = self.edit_canvas.canvasx(event.x), self.edit_canvas.canvasy(event.y)
        sx, sy = self.edit_data['drag_start']
        if self.edit_data.get('temp_shape'): self.edit_canvas.delete(self.edit_data['temp_shape'])
        c = self.edit_color_var.get()
        f = c if self.edit_fill_var.get() else ''
        t = self.edit_data['selected_tool']
        if t == 'rectangle': self.edit_data['temp_shape'] = self.edit_canvas.create_rectangle(sx, sy, cx, cy, outline=c, fill=f, dash=(4,4))
        elif t == 'circle': r=((cx-sx)**2+(cy-sy)**2)**0.5; self.edit_data['temp_shape'] = self.edit_canvas.create_oval(sx-r, sy-r, sx+r, sy+r, outline=c, fill=f, dash=(4,4))
        elif t == 'line': self.edit_data['temp_shape'] = self.edit_canvas.create_line(sx, sy, cx, cy, fill=c, dash=(4,4))

    def edit_canvas_release(self, event):
        if not self.edit_data.get('drag_start'): return
        cx, cy = self.edit_canvas.canvasx(event.x), self.edit_canvas.canvasy(event.y)
        sx, sy = self.edit_data['drag_start']
        if abs(cx-sx)<5 and abs(cy-sy)<5: self.edit_data['drag_start']=None; return
        scale = self.edit_data.get('scale_factor', 1.0)
        pdf_h = self.edit_data.get('pdf_height', 842)
        psx, psy = sx*scale, pdf_h-(sy*scale)
        pcx, pcy = cx*scale, pdf_h-(cy*scale)
        c = self.edit_color_var.get()
        t = self.edit_data['selected_tool']
        if t == 'rectangle':
            self.edit_data['edits'].append({'type':'rectangle','page':self.edit_data['current_page'],'params':{'x':min(psx,pcx),'y':max(psy,pcy),'width':abs(pcx-psx),'height':abs(psy-pcy),'fill':c if self.edit_fill_var.get() else 'none','stroke':c}})
        elif t == 'circle':
            r = ((cx-sx)**2+(cy-sy)**2)**0.5*scale
            self.edit_data['edits'].append({'type':'circle','page':self.edit_data['current_page'],'params':{'x':psx,'y':psy,'radius':r,'fill':c if self.edit_fill_var.get() else 'none','stroke':c}})
        elif t == 'line':
            self.edit_data['edits'].append({'type':'line','page':self.edit_data['current_page'],'params':{'x1':psx,'y1':psy,'x2':pcx,'y2':pcy,'color':c,'width':int(self.edit_size_var.get())}})
        self.edit_redraw_edits()
        self.edit_data['drag_start'] = None

    def edit_canvas_motion(self, event):
        if not self.edit_data['path'] or self.edit_data.get('drag_start'): return
        cx, cy = self.edit_canvas.canvasx(event.x), self.edit_canvas.canvasy(event.y)
        scale, pdf_h = self.edit_data.get('scale_factor', 1.0), self.edit_data.get('pdf_height', 842)
        self.edit_status.configure(text=f"X={int(cx*scale)}, Y={int(pdf_h-(cy*scale))}")

    def edit_add_text_at_position(self, cx, cy):
        scale, pdf_h = self.edit_data.get('scale_factor', 1.0), self.edit_data.get('pdf_height', 842)
        px, py = cx*scale, pdf_h-(cy*scale)
        
        dialog = ctk.CTkToplevel(self.root)
        dialog.title("Metin")
        dialog.geometry("300x150")
        dialog.transient(self.root)
        dialog.grab_set()
        
        ctk.CTkLabel(dialog, text=f"X={int(px)}, Y={int(py)}").pack(pady=5)
        e = ctk.CTkEntry(dialog, width=200)
        e.pack(pady=5)
        
        def save():
            if not e.get(): return
            self.edit_data['edits'].append({'type':'text','page':self.edit_data['current_page'],'params':{'text':e.get(),'x':px,'y':py,'size':int(self.edit_size_var.get()),'color':self.edit_color_var.get()}})
            self.edit_redraw_edits(); dialog.destroy()
            
        e.bind('<Return>', lambda ev: save())
        ctk.CTkButton(dialog, text="Ekle", command=save).pack(pady=5)

    def edit_prev_page(self):
        if self.edit_data['current_page'] > 1:
            self.edit_data['current_page'] -= 1
            self.edit_page_label.configure(text=f"Sayfa: {self.edit_data['current_page']} / {self.edit_data['total_pages']}")
            self.edit_refresh_preview()

    def edit_next_page(self):
        if self.edit_data['current_page'] < self.edit_data['total_pages']:
            self.edit_data['current_page'] += 1
            self.edit_page_label.configure(text=f"Sayfa: {self.edit_data['current_page']} / {self.edit_data['total_pages']}")
            self.edit_refresh_preview()

    def edit_add_image(self):
        path = filedialog.askopenfilename(filetypes=[("Resim", "*.png *.jpg *.jpeg")])
        if not path: return
        self.edit_data['edits'].append({'type':'image','page':self.edit_data['current_page'],'params':{'image_path':path,'x':50,'y':100,'width':200,'height':150}})
        self.edit_status.configure(text=f"Resim eklendi: {Path(path).name}")

    def edit_add_watermark(self):
        self.edit_data['edits'].append({'type':'watermark','page':self.edit_data['current_page'],'params':{'text':'GİZLİ','position':'center','size':48,'alpha':0.3}})
        self.edit_status.configure(text="Filigran eklendi")

    def edit_rotate_page(self):
        self.edit_data['edits'].append({'type':'rotate','page':self.edit_data['current_page'],'params':{'angle':90}})
        self.edit_status.configure(text="Döndürme eklendi")

    def edit_delete_page(self):
        if messagebox.askyesno("Onay", f"Sayfa {self.edit_data['current_page']} silinsin mi?"):
            self.edit_data['edits'].append({'type':'delete_page','page':self.edit_data['current_page'],'params':{}})
            self.edit_status.configure(text="Silme eklendi")

    def edit_show_edits(self):
        if not self.edit_data['edits']: return messagebox.showinfo("Bilgi", "Düzenleme yok.")
        d = ctk.CTkToplevel(self.root)
        d.title("Düzenlemeler")
        d.geometry("400x300")
        
        lb = tk.Listbox(d, bg="#2b2b2b" if ctk.get_appearance_mode()=="Dark" else "#ffffff", 
                        fg="white" if ctk.get_appearance_mode()=="Dark" else "black", highlightthickness=0)
        lb.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        for i, e in enumerate(self.edit_data['edits']): lb.insert(tk.END, f"{i+1}. {e['type']} - S:{e['page']}")
        
        def rem():
            if lb.curselection():
                del self.edit_data['edits'][lb.curselection()[0]]
                d.destroy(); messagebox.showinfo("Bilgi", "Silindi")
                
        ctk.CTkButton(d, text="Sil", command=rem, fg_color="#E74C3C", hover_color="#C0392B").pack(pady=10)

    def edit_clear_all(self):
        if messagebox.askyesno("Onay", "Tümü temizlensin mi?"):
            self.edit_data['edits'] = []
            self.edit_status.configure(text="Temizlendi")
            self.edit_redraw_edits()

    def edit_save_pdf(self):
        if not self.edit_data['path'] or not self.edit_data['edits']: return messagebox.showwarning("Uyarı", "Dosya/düzenleme eksik!")
        out = filedialog.asksaveasfilename(defaultextension=".pdf")
        if not out: return
        def run():
            self.root.after(0, lambda: self.edit_status.configure(text="Kaydediliyor..."))
            try:
                reader, writer = PdfReader(self.edit_data['path']), PdfWriter()
                del_pages = {e['page'] for e in self.edit_data['edits'] if e['type']=='delete_page'}
                for i in range(len(reader.pages)):
                    p_num = i + 1
                    if p_num in del_pages: continue
                    page = reader.pages[i]
                    p_edits = [e for e in self.edit_data['edits'] if e['page']==p_num and e['type']!='delete_page']
                    if p_edits:
                        ov = self._create_overlay_pdf(page, p_edits)
                        if ov: page.merge_page(PdfReader(ov).pages[0]); os.remove(ov)
                    for _ in [e for e in self.edit_data['edits'] if e['page']==p_num and e['type']=='rotate']: page.rotate(90)
                    writer.add_page(page)
                with open(out, 'wb') as f: writer.write(f)
                self.root.after(0, lambda: self.edit_status.configure(text="Kaydedildi"))
                messagebox.showinfo("Başarılı", out)
            except Exception as e: messagebox.showerror("Hata", str(e))
        threading.Thread(target=run, daemon=True).start()

    def _create_overlay_pdf(self, original_page, edits):
        try:
            fd, temp = tempfile.mkstemp(suffix=".pdf"); os.close(fd)
            w, h = float(original_page.mediabox.width), float(original_page.mediabox.height)
            c = canvas.Canvas(temp, pagesize=(w, h))
            cmap = {'black':rl_colors.black,'white':rl_colors.white,'red':rl_colors.red,'blue':rl_colors.blue,'green':rl_colors.green,'none':None}
            for e in edits:
                p = e['params']
                if e['type']=='text': c.setFillColor(cmap.get(p['color'],rl_colors.black)); c.setFont("Helvetica",p['size']); c.drawString(p['x'], h-p['y'], p['text'])
                elif e['type']=='rectangle':
                    c.rect(p['x'], h-p['y']-p['height'], p['width'], p['height'], fill=p['fill'] in ('white','black','red','blue','green'), stroke=p['stroke']!='none')
                elif e['type']=='circle':
                    c.circle(p['x'], h-p['y'], p['radius'], fill=p['fill'] in ('white','black','red','blue','green'), stroke=p['stroke']!='none')
                elif e['type']=='line': c.line(p['x1'], h-p['y1'], p['x2'], h-p['y2'])
                elif e['type']=='image' and os.path.exists(p['image_path']): c.drawImage(p['image_path'], p['x'], h-p['y']-p['height'], p['width'], p['height'], mask='auto')
                elif e['type']=='watermark':
                    c.saveState(); c.setFillGray(0.5, alpha=p['alpha']); c.setFont("Helvetica-Bold", p['size'])
                    if p['position']=='center': tw=c.stringWidth(p['text'],"Helvetica-Bold",p['size']); c.drawString((w-tw)/2, h/2, p['text'])
                    else: c.translate(w/2, h/2); c.rotate(45); tw=c.stringWidth(p['text'],"Helvetica-Bold",p['size']); c.drawString(-tw/2, 0, p['text'])
                    c.restoreState()
            c.save(); return temp
        except: return None

    # ================= PDF → DİĞER FORMATLAR =================
    def setup_export_ui(self):
        container = ctk.CTkFrame(self.export_tab, fg_color="transparent")
        container.pack(fill=tk.BOTH, expand=True, padx=30, pady=20)
        
        ctk.CTkLabel(container, text="📤 PDF'i Diğer Formatlara Dönüştür", font=ctk.CTkFont('Segoe UI', 18, 'bold')).pack(anchor='w')
        
        file_frame = ctk.CTkFrame(container)
        file_frame.pack(fill=tk.X, pady=(15, 10))
        self.export_file_label = ctk.CTkLabel(file_frame, text="Lütfen bir PDF dosyası seçin...", anchor='w', padx=15, font=('Segoe UI', 12))
        self.export_file_label.pack(side=tk.LEFT, fill=tk.X, expand=True, pady=10)
        ctk.CTkButton(file_frame, text="Dosya Seç", command=self.export_select_pdf, width=120, height=35).pack(side=tk.RIGHT, padx=10, pady=5)
        
        info_frame = ctk.CTkFrame(container, fg_color="transparent")
        info_frame.pack(fill=tk.X, pady=10)
        self.export_page_info = ctk.CTkLabel(info_frame, text="Seçili Sayfalar: -", font=('Segoe UI', 12))
        self.export_page_info.pack(side=tk.LEFT)
        ctk.CTkButton(info_frame, text="Sayfaları Düzenle", command=self.export_open_page_selector, fg_color="#9B59B6", hover_color="#8E44AD", width=120, height=35).pack(side=tk.RIGHT)
        
        settings_box = ctk.CTkFrame(container)
        settings_box.pack(fill=tk.X, pady=20)
        ctk.CTkLabel(settings_box, text=" Dönüşüm Ayarları ", font=ctk.CTkFont('Segoe UI', 14, 'bold')).pack(anchor="w", padx=15, pady=(15, 0))
        
        inner_set = ctk.CTkFrame(settings_box, fg_color="transparent")
        inner_set.pack(fill=tk.X, padx=15, pady=15)
        
        ctk.CTkLabel(inner_set, text="Hedef Format:", font=('Segoe UI', 12)).pack(side=tk.LEFT)
        self.export_format_var = ctk.StringVar(value="png")
        ctk.CTkComboBox(inner_set, variable=self.export_format_var, 
                     values=["PNG", "JPG", "DOCX", "DOC", "XLSX", "TXT", "CSV", "PPTX"], state="readonly", width=120).pack(side=tk.LEFT, padx=10)
                     
        ctk.CTkLabel(inner_set, text="DPI:", font=('Segoe UI', 12)).pack(side=tk.LEFT, padx=(20,0))
        self.export_dpi_var = ctk.StringVar(value="150")
        ctk.CTkComboBox(inner_set, variable=self.export_dpi_var, values=["72", "150", "300", "600"], state="readonly", width=100).pack(side=tk.LEFT, padx=10)
        
        self.export_progress = ctk.CTkProgressBar(container)
        self.export_progress.set(0)
        self.export_progress.pack(fill=tk.X, pady=15)
        self.export_status = ctk.CTkLabel(container, text="Hazır.", font=('Segoe UI', 11), text_color="gray", anchor='w')
        self.export_status.pack(anchor='w', pady=(0, 10))
        
        ctk.CTkButton(container, text="🚀 DÖNÜŞTÜR", command=self.start_export, fg_color="#27AE60", hover_color="#229954", font=('Segoe UI', 16, 'bold'), width=200, height=40).pack(pady=5)

    def export_select_pdf(self):
        path = filedialog.askopenfilename(filetypes=[("PDF", "*.pdf")])
        if path:
            try:
                r = PdfReader(path)
                self.export_data = {'path': path, 'name': Path(path).name, 'total_pages': len(r.pages), 'selected_pages': list(range(1, len(r.pages) + 1))}
                self.export_file_label.configure(text=self.export_data['name'])
                self._update_export_ui()
            except Exception as e: messagebox.showerror("Hata", str(e))

    def export_open_page_selector(self):
        if not self.export_data['path']: return messagebox.showwarning("Uyarı", "Dosya seçin!")
        PageSelectorDialog(self.root, self.export_data, self._update_export_ui)

    def _update_export_ui(self):
        s = self.export_data['selected_pages']
        if not s: self.export_page_info.configure(text="-"); return
        ranges, start, end = [], s[0], s[0]
        for i in range(1, len(s)):
            if s[i]==end+1: end=s[i]
            else: ranges.append(f"{start}-{end}" if start!=end else str(start)); start=end=s[i]
        ranges.append(f"{start}-{end}" if start!=end else str(start))
        self.export_page_info.configure(text=f"Seçili ({len(s)}): {', '.join(ranges)}")

    def start_export(self):
        if not self.export_data['path']: return messagebox.showwarning("Uyarı", "Dosya seçilmedi!")
        if not self.export_data['selected_pages']: return messagebox.showwarning("Uyarı", "Sayfa seçin!")
        out_dir = filedialog.askdirectory()
        if not out_dir: return
        raw = self.export_format_var.get().strip().lower()
        if raw.startswith("jpg"): fmt = "jpg"
        elif raw.startswith("png"): fmt = "png"
        elif "docx" in raw: fmt = "docx"
        elif "doc" in raw: fmt = "doc"
        elif "xlsx" in raw: fmt = "xlsx"
        elif "txt" in raw: fmt = "txt"
        elif "csv" in raw: fmt = "csv"
        elif "pptx" in raw: fmt = "pptx"
        else: fmt = raw
        self.root.after(0, lambda: self.export_status.configure(text="İşlem başlatılıyor..."))
        threading.Thread(target=self._run_export, args=(out_dir, fmt, int(self.export_dpi_var.get())), daemon=True).start()

    def _run_export(self, out_dir, fmt, dpi):
        total = len(self.export_data['selected_pages'])
        self.root.after(0, lambda: self.export_progress.set(0))
        try:
            if fmt in ('png', 'jpg'): res = self._export_images(out_dir, fmt, dpi, Path(self.export_data['path']).stem)
            elif fmt in ('doc', 'docx'): res = self._export_to_docx(out_dir, fmt, Path(self.export_data['path']).stem)
            elif fmt == 'xlsx': res = self._export_to_excel(out_dir, Path(self.export_data['path']).stem)
            elif fmt == 'txt': res = self._export_txt(out_dir, Path(self.export_data['path']).stem)
            elif fmt == 'csv': res = self._export_csv(out_dir, Path(self.export_data['path']).stem)
            elif fmt == 'pptx': res = self._export_pptx(out_dir, Path(self.export_data['path']).stem)
            else: raise ValueError("Desteklenmeyen format!")
            self.root.after(0, lambda: self.export_status.configure(text=f"✅ Tamamlandı. {res} sayfa işlendi."))
            messagebox.showinfo("Başarılı", f"Dönüşüm bitti.\nKlasör: {out_dir}")
        except Exception as e:
            self.root.after(0, lambda: self.export_status.configure(text=f"❌ Hata: {e}"))
            messagebox.showerror("Hata", str(e))

    def _export_images(self, out_dir, fmt, dpi, base):
        if not PYPDFIUM_AVAILABLE: raise RuntimeError("pip install pypdfium2")
        pdf = pdfium.PdfDocument(self.export_data['path'])
        total = len(self.export_data['selected_pages'])
        for i, p in enumerate(sorted(self.export_data['selected_pages'])):
            self.root.after(0, lambda text=f"Render: {p}": self.export_status.configure(text=text))
            img = pdf[p-1].render(scale=dpi/72.0).to_pil()
            img.save(os.path.join(out_dir, f"{base}_p{p}.{fmt}"), fmt.upper())
            self.root.after(0, lambda v=(i+1)/total: self.export_progress.set(v))
        return i+1

    def _export_to_docx(self, out_dir, fmt, base):
        tmp = os.path.join(tempfile.gettempdir(), f"tmp_{os.getpid()}.pdf")
        try:
            w, r = PdfWriter(), PdfReader(self.export_data['path'])
            for p in sorted(self.export_data['selected_pages']): w.add_page(r.pages[p-1])
            with open(tmp, "wb") as f: w.write(f)
            out = os.path.join(out_dir, f"{base}.{fmt}")
            if PDF2DOCX_AVAILABLE:
                try:
                    self.root.after(0, lambda: self.export_status.configure(text="pdf2docx dönüştürüyor..."))
                    Converter(tmp).convert(out); return 1
                except: pass
            self.root.after(0, lambda: self.export_status.configure(text="LibreOffice deneniyor..."))
            return self._convert_with_libreoffice(tmp, out_dir, base, fmt)
        finally:
            if os.path.exists(tmp): os.remove(tmp)

    def _export_to_excel(self, out_dir, base):
        tmp = os.path.join(tempfile.gettempdir(), f"tmp_{os.getpid()}.pdf")
        try:
            w, r = PdfWriter(), PdfReader(self.export_data['path'])
            for p in sorted(self.export_data['selected_pages']): w.add_page(r.pages[p-1])
            with open(tmp, "wb") as f: w.write(f)
            return self._convert_with_libreoffice(tmp, out_dir, base, "xlsx")
        finally:
            if os.path.exists(tmp): os.remove(tmp)

    def _export_txt(self, out_dir, base):
        r = PdfReader(self.export_data['path'])
        text = []
        for p in sorted(self.export_data['selected_pages']):
            t = r.pages[p-1].extract_text()
            if t: text.append(t)
        with open(os.path.join(out_dir, f"{base}.txt"), "w", encoding="utf-8") as f: f.write("\n\n".join(text))
        return 1

    def _export_csv(self, out_dir, base):
        r = PdfReader(self.export_data['path'])
        rows = []
        for p in sorted(self.export_data['selected_pages']):
            t = r.pages[p-1].extract_text()
            if t: rows.extend([line.strip() for line in t.splitlines() if line.strip()])
        with open(os.path.join(out_dir, f"{base}.csv"), "w", encoding="utf-8", newline="") as f:
            csv.writer(f).writerows([[row] for row in rows])
        return 1

    def _export_pptx(self, out_dir, base):
        if not PPTX_AVAILABLE: raise RuntimeError("pip install python-pptx")
        if not PYPDFIUM_AVAILABLE: raise RuntimeError("pip install pypdfium2")
        prs = Presentation()
        slide_w, slide_h = Inches(11.69), Inches(8.27)
        prs.slide_width, prs.slide_height = slide_w, slide_h
        pdf = pdfium.PdfDocument(self.export_data['path'])
        total = len(self.export_data['selected_pages'])
        for i, p in enumerate(sorted(self.export_data['selected_pages'])):
            self.root.after(0, lambda text=f"Sunu: {p}": self.export_status.configure(text=text))
            img_path = os.path.join(out_dir, f"_temp_{p}.png")
            pdf[p-1].render(scale=2.0).to_pil().save(img_path)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), slide_w, slide_h)
            if os.path.exists(img_path): os.remove(img_path)
            self.root.after(0, lambda v=(i+1)/total: self.export_progress.set(v))
        prs.save(os.path.join(out_dir, f"{base}.pptx"))
        return i+1

    def _convert_with_libreoffice(self, inp, out_dir, base, ext):
        lo = shutil.which('libreoffice') or shutil.which('soffice') or r'C:\Program Files\LibreOffice\program\soffice.exe'
        if not os.path.exists(lo): raise RuntimeError("LibreOffice kurulu değil!")
        prof = os.path.join(tempfile.gettempdir(), f"lo_{os.getpid()}")
        os.makedirs(prof, exist_ok=True)
        res = subprocess.run([lo, '--headless', '--norestore', f'-env:UserInstallation={Path(prof).as_uri()}', '--convert-to', ext, '--outdir', out_dir, inp], capture_output=True, timeout=90)
        if res.returncode != 0: raise RuntimeError(f"LO Hatası: {res.stderr.decode('utf-8', errors='replace')}")
        return 1


if __name__ == "__main__":
    root = ctk.CTk()
    app = PDFToolboxApp(root)
    root.mainloop()
